import os
import io
import re
import csv
import time
from typing import List, Tuple

import streamlit as st
import pandas as pd

from pypdf import PdfReader
import docx2txt
from pptx import Presentation

# OpenAI SDK
try:
    from openai import OpenAI
    _HAS_OPENAI = True
except Exception:
    _HAS_OPENAI = False

# ─────────────────────────────────────
# Secrets / Env
# ─────────────────────────────────────
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", st.secrets.get("OPENAI_API_KEY", ""))
OPENAI_MODEL = os.getenv("OPENAI_MODEL", st.secrets.get("OPENAI_MODEL", "gpt-4o-mini"))
OPENAI_PROJECT = os.getenv("OPENAI_PROJECT", st.secrets.get("OPENAI_PROJECT", ""))  # optionnel

# ─────────────────────────────────────
# UI – Sidebar
# ─────────────────────────────────────
st.set_page_config(page_title="IA Révisions Étudiant – MVP", page_icon="📚", layout="wide")
st.sidebar.title("⚙️ Configuration")
model = st.sidebar.text_input("Modèle (OpenAI)", value=OPENAI_MODEL)
max_flashcards = st.sidebar.slider("Nombre de flashcards", 5, 100, 20)
max_qcm = st.sidebar.slider("Nombre de QCM", 5, 50, 10)
summary_style = st.sidebar.selectbox(
    "Style de fiche",
    ["Ultra-court (post-it)", "Standard (fiche bac)", "Détaillé (type prof)"]
)
target_lang_label = st.sidebar.selectbox(
    "Traduire la fiche/flashcards/QCM dans :", ["(aucune)", "fr", "en", "nl", "de", "es", "it"], index=0
)

# progression (simple compteur d’actions réussies)
if "prog_total" not in st.session_state:
    st.session_state.prog_total = 0
if "prog_done" not in st.session_state:
    st.session_state.prog_done = 0

st.title("📚 IA Révisions Étudiant – MVP")
st.caption("Upload ton cours → Fiches + Flashcards + QCM + Chat + Traductions.")

# ─────────────────────────────────────
# File parsers
# ─────────────────────────────────────
def extract_text_from_pdf(file) -> str:
    reader = PdfReader(file)
    texts = []
    for page in reader.pages:
        try:
            texts.append(page.extract_text() or "")
        except Exception:
            pass
    return "\n".join(texts)

def extract_text_from_docx(file) -> str:
    data = file.read()
    tmp_path = "_tmp_upload.docx"
    with open(tmp_path, "wb") as f:
        f.write(data)
    text = docx2txt.process(tmp_path) or ""
    os.remove(tmp_path)
    return text

def extract_text_from_pptx(file) -> str:
    prs = Presentation(file)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

ALLOWED_TYPES = {
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "pptx": extract_text_from_pptx,
    "txt": lambda f: f.read().decode("utf-8", errors="ignore"),
}

# ─────────────────────────────────────
# Prompts
# ─────────────────────────────────────
PROMPT_SUMMARY = (
    "Tu es un tuteur pédagogique. Fais une fiche de révision à partir du texte suivant. "
    "Propose 3 sections: 1) Points clés (bullet points), 2) Définitions essentielles, 3) Exemples/Application. "
    "Adapte la longueur au style demandé: {style}. \n\nTEXTE:\n{content}\n\nRéponds en Markdown propre."
)

PROMPT_FLASHCARDS = (
    "Génère des paires de flashcards (Question;Réponse) à partir du texte suivant. "
    "Questions courtes, réponses précises. Donne {n} paires. Sépare chaque paire sur une ligne via le séparateur ' | '. "
    "TEXTE:\n{content}"
)

PROMPT_QCM = (
    "Génère {n} QCM à partir du texte suivant. Pour chaque QCM, propose:\n"
    "- Énoncé \n- 4 options (A,B,C,D) \n- Indique la bonne réponse (A-D) et une brève justification.\n"
    "Format:\nQ: ...\nA) ...\nB) ...\nC) ...\nD) ...\nRéponse: X | Justification: ...\n---\n"
    "TEXTE:\n{content}"
)

# ─────────────────────────────────────
# OpenAI call with retry/backoff (gère 429)
# ─────────────────────────────────────
def call_openai(prompt: str, temperature: float = 0.2) -> str:
    if not _HAS_OPENAI:
        raise RuntimeError("Le package openai n'est pas installé.")
    if not OPENAI_API_KEY:
        raise RuntimeError("OPENAI_API_KEY absent. Ajoute-le dans Secrets.")
    client = OpenAI(api_key=OPENAI_API_KEY, project=OPENAI_PROJECT or None)

    delays = [0, 1, 2, 4, 8]  # secondes
    last_err = None
    for d in delays:
        if d:
            time.sleep(d)
        try:
            resp = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": "Tu es un expert de la pédagogie claire."},
                    {"role": "user", "content": prompt},
                ],
                temperature=temperature,
            )
            return resp.choices[0].message.content.strip()
        except Exception as e:
            err = str(e)
            last_err = e
            if "429" in err or "rate" in err.lower() or "temporar" in err.lower():
                continue  # on retente
            raise RuntimeError(f"Erreur OpenAI: {e}")
    raise RuntimeError(f"Erreur OpenAI (après plusieurs tentatives): {last_err}")

# ─────────────────────────────────────
# Utils (normalisation, exports)
# ─────────────────────────────────────
def normalize_text(s: str) -> str:
    return re.sub(r"\s+", " ", s or " ").strip()

def to_anki_csv(pairs: List[Tuple[str, str]]) -> bytes:
    buf = io.StringIO()
    writer = csv.writer(buf)
    for q, a in pairs:
        writer.writerow([q, a])
    return buf.getvalue().encode("utf-8")

def to_quizlet_tsv(pairs: List[Tuple[str, str]]) -> bytes:
    # Quizlet: Term<TAB>Definition par ligne
    lines = [f"{q}\t{a}" for q, a in pairs]
    return ("\n".join(lines)).encode("utf-8")

# ─────────── Traductions ───────────
def translate_text(text: str, target_lang: str) -> str:
    """Traduit un bloc de texte en conservant la structure Markdown."""
    prompt = (
        f"Traduis le texte suivant en langue code '{target_lang}'. "
        f"Conserve EXACTEMENT la structure et la mise en forme (titres, listes, lignes). "
        f"Ne rajoute pas de commentaire.\n\nTEXTE:\n{text}"
    )
    return call_openai(prompt, temperature=0.0)

def translate_flashcards(pairs: List[Tuple[str, str]], target_lang: str) -> List[Tuple[str, str]]:
    """Traduit une liste de flashcards en gardant le format 'Question | Réponse' par ligne."""
    if not pairs:
        return []
    joined = "\n".join([f"{q} | {a}" for q, a in pairs])
    prompt = (
        f"Tu es traducteur. Traduis CHAQUE paire suivante en langue code '{target_lang}'. "
        f"IMPORTANT: Garde EXACTEMENT le même format 'Question | Réponse' par ligne, "
        f"le même nombre de lignes et le même ordre. Réponds UNIQUEMENT avec les lignes traduites.\n\n"
        f"{joined}"
    )
    out = call_openai(prompt, temperature=0.0)
    translated_pairs = []
    for line in out.splitlines():
        if "|" in line:
            q, a = [t.strip(" -:\t") for t in line.split("|", 1)]
            if q and a:
                translated_pairs.append((q, a))
    # si l'IA a renvoyé moins de lignes, on garde au moins l'original
    return translated_pairs if len(translated_pairs) == len(pairs) else pairs

def translate_qcm(qcm_text: str, target_lang: str) -> str:
    """Traduit un QCM en gardant EXACTEMENT la structure attendue."""
    prompt = (
        f"Traduis le QCM ci-dessous en langue code '{target_lang}' en conservant EXACTEMENT cette structure :\n"
        f"Les balises 'Q:', 'A)', 'B)', 'C)', 'D)', 'Réponse:', 'Justification:' et les séparateurs '---'.\n"
        f"N'ajoute rien, ne reformate pas, ne supprime pas de questions.\n\n"
        f"{qcm_text}"
    )
    return call_openai(prompt, temperature=0.0)

# ─────────────────────────────────────
# Main
# ─────────────────────────────────────
uploaded = st.file_uploader(
    "Téléverse tes cours (PDF/DOCX/PPTX/TXT)",
    type=list(ALLOWED_TYPES.keys()),
    accept_multiple_files=True
)

if uploaded:
    st.success(f"{len(uploaded)} fichier(s) uploadé(s).")

    all_texts = []
    for file in uploaded:
        ext = file.name.split(".")[-1].lower()
        parser = ALLOWED_TYPES.get(ext)
        if not parser:
            st.warning(f"Type non supporté: {file.name}")
            continue
        with st.spinner(f"Extraction du texte: {file.name}"):
            try:
                text = parser(file)
                text = normalize_text(text)
            except Exception as e:
                st.error(f"Erreur extraction {file.name}: {e}")
                text = ""
        if text:
            st.write(f"**Extrait ({file.name}) – 1ère page env.**")
            st.text(text[:800] + ("..." if len(text) > 800 else ""))
            all_texts.append(text)

    merged_text = "\n\n".join(all_texts)

    # Limite la taille envoyée au modèle pour éviter 429 et coûts
    MAX_CHARS = 8000
    if len(merged_text) > MAX_CHARS:
        st.info(f"Le document est long ({len(merged_text)} caractères). On traite d'abord ~{MAX_CHARS} caractères.")
        merged_text = merged_text[:MAX_CHARS]

    if merged_text.strip():
        # Progression
        with st.expander("📈 Progression"):
            total = st.session_state.prog_total or 1
            ratio = (st.session_state.prog_done / total) if total else 0
            st.progress(min(ratio, 1.0), text=f"{int(ratio*100)} % des actions complétées")

        col1, col2, col3 = st.columns(3)

        # ── FICHES ─────────────────────
        with col1:
            if st.button("📝 Générer Fiche(s)"):
                with st.spinner("Génération des fiches..."):
                    prompt = PROMPT_SUMMARY.format(style=summary_style, content=merged_text)
                    summary_md = call_openai(prompt)
                st.subheader("Fiche de révision")
                st.markdown(summary_md)
                st.download_button("Télécharger la fiche (MD)", summary_md.encode("utf-8"), file_name="fiche_revision.md")

                # Traduction optionnelle de la fiche
                if target_lang_label != "(aucune)":
                    with st.spinner(f"Traduction de la fiche en {target_lang_label}..."):
                        summary_tr = translate_text(summary_md, target_lang_label)
                    st.subheader(f"Fiche (traduite : {target_lang_label})")
                    st.markdown(summary_tr)
                    st.download_button(
                        f"Télécharger la fiche traduite (MD, {target_lang_label})",
                        summary_tr.encode("utf-8"),
                        file_name=f"fiche_revision_{target_lang_label}.md"
                    )

                st.session_state.prog_total += 1
                st.session_state.prog_done += 1

        # ── FLASHCARDS ─────────────────
        with col2:
            if st.button("🃏 Générer Flashcards"):
                with st.spinner("Génération des flashcards..."):
                    prompt = PROMPT_FLASHCARDS.format(n=max_flashcards, content=merged_text)
                    raw = call_openai(prompt)
                pairs = []
                for line in raw.splitlines():
                    if "|" in line:
                        q, a = [t.strip(" -:") for t in line.split("|", 1)]
                        if q and a:
                            pairs.append((q, a))
                if not pairs:
                    st.warning("Aucune paire détectée. Vérifie le contenu du cours.")
                # tableau original
                df = pd.DataFrame(pairs, columns=["Front (Question)", "Back (Réponse)"])
                st.dataframe(df, use_container_width=True)
                st.download_button(
                    "Télécharger Anki CSV",
                    to_anki_csv(pairs),
                    file_name="flashcards_anki.csv",
                    mime="text/csv"
                )
                st.download_button(
                    "Télécharger Quizlet TSV",
                    to_quizlet_tsv(pairs),
                    file_name="flashcards_quizlet.tsv",
                    mime="text/tab-separated-values"
                )

                # Traduction optionnelle des flashcards
                if pairs and target_lang_label != "(aucune)":
                    with st.spinner(f"Traduction des flashcards en {target_lang_label}..."):
                        pairs_tr = translate_flashcards(pairs, target_lang_label)
                    df_tr = pd.DataFrame(pairs_tr, columns=[f"Front ({target_lang_label})", f"Back ({target_lang_label})"])
                    st.dataframe(df_tr, use_container_width=True)
                    st.download_button(
                        f"Télécharger Anki CSV ({target_lang_label})",
                        to_anki_csv(pairs_tr),
                        file_name=f"flashcards_anki_{target_lang_label}.csv",
                        mime="text/csv"
                    )
                    st.download_button(
                        f"Télécharger Quizlet TSV ({target_lang_label})",
                        to_quizlet_tsv(pairs_tr),
                        file_name=f"flashcards_quizlet_{target_lang_label}.tsv",
                        mime="text/tab-separated-values"
                    )

                st.session_state.prog_total += 1
                st.session_state.prog_done += 1

        # ── QCM ────────────────────────
        with col3:
            if st.button("✅ Générer QCM"):
                with st.spinner("Génération des QCM..."):
                    prompt = PROMPT_QCM.format(n=max_qcm, content=merged_text)
                    qcm_text = call_openai(prompt)
                st.subheader("QCM")
                st.text(qcm_text)
                st.download_button("Télécharger QCM (TXT)", qcm_text.encode("utf-8"), file_name="qcm.txt")

                # Traduction optionnelle du QCM
                if target_lang_label != "(aucune)":
                    with st.spinner(f"Traduction du QCM en {target_lang_label}..."):
                        qcm_tr = translate_qcm(qcm_text, target_lang_label)
                    st.subheader(f"QCM (traduit : {target_lang_label})")
                    st.text(qcm_tr)
                    st.download_button(
                        f"Télécharger QCM traduit (TXT, {target_lang_label})",
                        qcm_tr.encode("utf-8"),
                        file_name=f"qcm_{target_lang_label}.txt"
                    )

                st.session_state.prog_total += 1
                st.session_state.prog_done += 1

        # ── CHAT AVEC LE COURS ────────
        st.markdown("---")
        st.subheader("💬 Chat avec ton cours")
        st.caption("Pose une question sur le contenu téléversé. L'IA répond en se basant uniquement sur ce texte.")
        user_q = st.text_input("Ta question (ex: Explique la période réfractaire en 2 phrases)")
        if user_q and merged_text.strip():
            if st.button("Poser la question"):
                with st.spinner("Réflexion..."):
                    chat_prompt = (
                        "Tu es un tuteur. Réponds de manière courte et claire "
                        "en te basant EXCLUSIVEMENT sur le CONTEXTE ci-dessous. "
                        "Si l'info n'est pas dans le contexte, dis-le.\n\n"
                        f"CONTEXTE:\n{merged_text}\n\nQUESTION:\n{user_q}"
                    )
                    answer = call_openai(chat_prompt, temperature=0.0)
                st.markdown("**Réponse :**")
                st.write(answer)

    else:
        st.info("Upload un fichier pour commencer.")

# Footer
st.markdown("---")
st.caption("MVP pédagogique – Streamlit + OpenAI. Exports Anki/Quizlet, Chat, Traductions, Progression.")

